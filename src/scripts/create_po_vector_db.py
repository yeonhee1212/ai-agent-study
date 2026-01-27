import base64
import sys
from pathlib import Path

# 프로젝트 루트를 sys.path에 추가 (스크립트 직접 실행 시)
project_root = Path(__file__).parent.parent.parent
sys.path.insert(0, str(project_root))

from langchain_core.messages import HumanMessage
from pptx import Presentation
from pdf2image import convert_from_path
from src.agents.common import llm_model
from langchain_core.documents import Document
from langchain_chroma import Chroma
from langchain_huggingface import HuggingFaceEmbeddings

# =========================
# PPT
#  ├─ extract_ppt_text()  → 노트/텍스트
#  ├─ pdf_to_images()
#  │       └─ Vision 요약
#  └─ slide 단위로 병합 → LangChain Document
# =========================

PPT_PATH = Path("doc/ProObject21.pptx")
PDF_PATH = Path("doc/ProObject21.pdf")
IMG_OUT_DIR = Path("doc/images")

# =========================
# 1. PPT 텍스트 + 노트 추출
# =========================

def extract_ppt_text(ppt_path: Path):
    prs = Presentation(str(ppt_path))
    slides = []

    for i, slide in enumerate(prs.slides, start=1):
        title = slide.shapes.title.text if slide.shapes.title else ""

        body_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text.strip()
                if text:
                    body_text.append(text)

        notes = ""
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        slides.append({
            "slide": i,
            "title": title,
            "body_text": body_text,
            "notes": notes
        })

    return slides


def pdf_to_images(pdf_path: Path, out_dir: Path):
    out_dir.mkdir(parents=True, exist_ok=True)

    images = convert_from_path(
        str(pdf_path),  # 첫 번째 인자는 위치 인자로 경로 문자열 전달
        dpi=200,        # Vision엔 150~200이면 충분
        fmt="png"
    )

    image_paths = []
    for i, img in enumerate(images, start=1):
        img_path = out_dir / f"slide_{i}.png"
        img.save(img_path)
        image_paths.append(img_path)

    return image_paths

# 이미지 요약
def vision_summarize(pdf_path: Path) -> str:
    image_base64 = base64.b64encode(pdf_path.read_bytes()).decode()

    prompt = """
        This image is a system design PPT slide.
        Summarize it with a focus on:

        1. The main concepts or components

        2. The relationships represented by boxes, arrows, and hierarchical structures

        3. The overall flow illustrated by the diagram

        Write in English,
        and present only the key points without unnecessary introduction.
    """
    message = HumanMessage(
        content=[
            {"type": "text", "text": prompt},
            {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{image_base64}"}}
        ],
        max_tokens=100
    )
    response = llm_model.invoke([message])

    return str(response.content)

def merge_slide_data(
    ppt_slides: list[dict],
    image_paths: list[Path],
    vision_results: dict[int, str]) -> list[dict] | None:
    """
    slide 번호 기준으로
    - PPT 텍스트 / 노트
    - Vision 요약
    을 하나의 LangChain Document로 병합
    """

    merged = []

    for slide in ppt_slides:
        slide_no = slide["slide"]
        vision_text = vision_results.get(slide_no, "")

        title = slide.get("title", "")
        body_text = slide.get("body_text", [])
        notes = slide.get("notes", "")

        page_content = f"""
[슬라이드 번호]
{slide_no}

[제목]
{title}

[슬라이드 본문 텍스트]
{chr(10).join(body_text)}

[슬라이드 노트]
{notes}

[슬라이드 시각적 구조 요약]
{vision_text}
        """.strip()

        merged.append({
            "slide": slide_no,
            "content": page_content
        })
        
    return merged

def to_documents(merged_slides: list[dict]):
    documents = []

    for slide in merged_slides:
        documents.append(
            Document(
                page_content=slide["content"],
                metadata={
                    "source": "ppt",
                    "slide": slide["slide"]
                }
            )
        )

    return documents


slides = extract_ppt_text(PPT_PATH)
image_paths = pdf_to_images(PDF_PATH, IMG_OUT_DIR)

vision_results = {}
# 처음 4개 슬라이드만 처리
for i, img_path in enumerate(image_paths, start=1):
    print(f"Vision 요약 중: slide {i}")
    vision_results[i] = vision_summarize(img_path)

merged_slides = merge_slide_data(slides, image_paths, vision_results)

document_list = to_documents(merged_slides or [])

# 임베딩
embedding_model = HuggingFaceEmbeddings(
    model_name='jhgan/ko-sroberta-nli',
    model_kwargs={'device':'cpu'},
    encode_kwargs={'normalize_embeddings':True},
)
# 벡터 스토어에 저장
vector_store = Chroma.from_documents(
    documents=document_list,
    embedding=embedding_model,
    collection_name = 'rag_collection',
    persist_directory = './rag_collection'
)

retriever = vector_store.as_retriever(search_kwargs={'k': 3})

# print(retriever.invoke("framework란 뭐야?"))

# print(retriever.invoke("ProObject21의 주요기능은 크게 뭘로 나눠져있어?"))